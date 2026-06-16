from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ── Color Palette (4Persona Quiet Luxury) ──
BG_DARK = RGBColor(0x11, 0x11, 0x10)
BG_WARM = RGBColor(0xf8, 0xf6, 0xf2)
GOLD = RGBColor(0xc9, 0xa9, 0x6e)
GOLD_DARK = RGBColor(0xa8, 0x88, 0x4d)
TEXT = RGBColor(0x1c, 0x1b, 0x19)
TEXT_MUTED = RGBColor(0x6b, 0x65, 0x60)
WHITE = RGBColor(0xff, 0xff, 0xff)
API_RED = RGBColor(0xb8, 0x32, 0x32)
AIR_BLUE = RGBColor(0x3b, 0x77, 0xbc)
ANGIN_GOLD = RGBColor(0xc5, 0x96, 0x3a)
TANAH_GREEN = RGBColor(0x3a, 0x5c, 0x3a)

def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape_bg(slide, left, top, width, height, color, opacity=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_text_box(slide, left, top, width, height, text, font_size=18, color=TEXT, bold=False, alignment=PP_ALIGN.LEFT, font_name='Calibri'):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_bullet_slide(slide, left, top, width, height, items, font_size=16, color=TEXT, font_name='Calibri'):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = font_name
        p.space_after = Pt(8)
    return txBox

def add_accent_line(slide, left, top, width, color=GOLD):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Pt(3))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()

# ═══════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_slide_bg(slide, BG_DARK)

add_shape_bg(slide, Inches(0), Inches(0), Inches(0.15), Inches(7.5), GOLD)
add_text_box(slide, Inches(1.5), Inches(2.0), Inches(10), Inches(1.5),
    '4Persona Vundiego', font_size=54, color=GOLD, bold=True, alignment=PP_ALIGN.LEFT)
add_text_box(slide, Inches(1.5), Inches(3.5), Inches(10), Inches(1),
    'Personality Test Platform — Quiet Luxury Perfume Brand', font_size=26, color=WHITE, alignment=PP_ALIGN.LEFT)
add_accent_line(slide, Inches(1.5), Inches(4.5), Inches(3), GOLD)
add_text_box(slide, Inches(1.5), Inches(5.0), Inches(10), Inches(1),
    'Lead Generation & Sales Conversion Tool for Vun Diego (Maison de Parfum)', font_size=16, color=RGBColor(0xa0, 0x98, 0x90), alignment=PP_ALIGN.LEFT)
add_text_box(slide, Inches(1.5), Inches(6.2), Inches(10), Inches(0.5),
    'React 19 · NestJS 11 · PostgreSQL · Supabase · Prisma 6', font_size=14, color=GOLD, alignment=PP_ALIGN.LEFT)

# ═══════════════════════════════════════════════════════════════
# SLIDE 2 — WHAT IS THIS PROJECT?
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_WARM)
add_accent_line(slide, Inches(0.8), Inches(0.8), Inches(2), GOLD)
add_text_box(slide, Inches(0.8), Inches(1.0), Inches(11), Inches(0.8),
    'Project Overview', font_size=38, color=TEXT, bold=True)
add_text_box(slide, Inches(0.8), Inches(2.0), Inches(11), Inches(0.5),
    'A personality test web app for the perfume brand Vun Diego', font_size=18, color=TEXT_MUTED)

items = [
    'Users answer 25 questions → classified into 1 of 4 elements (Fire, Water, Wind, Earth)',
    'Each element maps to a Vundiego perfume: Choleric, Melancholic, Sanguine, Phlegmatic',
    'Core Purpose: Collect customer emails (leads) & drive perfume purchases',
    'Target Audience: Indonesian-speaking users (Bahasa Indonesia UI)',
    'Built as a lead generation & sales conversion tool',
    'User journey: Home → Auth → Test → Teaser (reveal) → Report (with shop links)',
]
add_bullet_slide(slide, Inches(0.8), Inches(2.8), Inches(11), Inches(4.5), items, 16, TEXT)

# ═══════════════════════════════════════════════════════════════
# SLIDE 3 — TECH STACK
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_accent_line(slide, Inches(0.8), Inches(0.8), Inches(2), GOLD)
add_text_box(slide, Inches(0.8), Inches(1.0), Inches(11), Inches(0.8),
    'Technology Stack', font_size=38, color=WHITE, bold=True)

stack_items = [
    ('Frontend', 'React 19 + Vite 8, Framer Motion 12, Recharts 3, Three.js / react-three-fiber'),
    ('Backend', 'NestJS 11 (TypeScript), modular architecture, JWT + Passport auth'),
    ('Database', 'PostgreSQL hosted on Supabase, Prisma 6 ORM with migrations'),
    ('Auth', 'JWT (email/password) + Google OAuth 2.0, bcrypt hashing (12 rounds)'),
    ('Email', 'Brevo (Sendinblue) REST API v3 — 300 emails/day free tier'),
    ('Storage', 'Supabase Storage for admin-uploaded report backgrounds'),
    ('Charts', 'Recharts radar chart for 4-element persona scores'),
    ('Hosting', 'Vercel (frontend) · Railway / Render (backend) — planned'),
]
y = 2.2
for label, desc in stack_items:
    add_text_box(slide, Inches(0.8), Inches(y), Inches(2.8), Inches(0.4),
        label, font_size=16, color=GOLD, bold=True)
    add_text_box(slide, Inches(3.8), Inches(y), Inches(8.5), Inches(0.4),
        desc, font_size=14, color=RGBColor(0xa0, 0x98, 0x90))
    y += 0.6

# ═══════════════════════════════════════════════════════════════
# SLIDE 4 — ARCHITECTURE
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_WARM)
add_accent_line(slide, Inches(0.8), Inches(0.8), Inches(2), GOLD)
add_text_box(slide, Inches(0.8), Inches(1.0), Inches(11), Inches(0.8),
    'System Architecture', font_size=38, color=TEXT, bold=True)

arch_items = [
    'Frontend (React SPA) ←→ Backend API (NestJS) ←→ PostgreSQL (Supabase)',
    'Axios HTTP client with JWT interceptor for authenticated requests',
    'Public report pages via token-based access (UUID, 30-day expiry)',
    'Visitor counting with sessionStorage de-duplication (1x per 24h)',
    'Homepage stats cached in-memory (5-min TTL) to reduce DB load',
    'Email sent asynchronously via Brevo after test submission',
    'Social share: html2canvas generates persona cards client-side',
]
add_bullet_slide(slide, Inches(0.8), Inches(2.2), Inches(11), Inches(4.5), arch_items, 16, TEXT)

# Architecture diagram description
add_shape_bg(slide, Inches(0.8), Inches(4.8), Inches(11.5), Inches(2.2), RGBColor(0xef, 0xec, 0xe6))
add_text_box(slide, Inches(1.2), Inches(4.9), Inches(3), Inches(0.4),
    'User Browser', font_size=14, color=TEXT, bold=True, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(1.2), Inches(5.2), Inches(3), Inches(0.3),
    'React · Vite · Framer Motion', font_size=11, color=TEXT_MUTED, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(5.0), Inches(4.9), Inches(3), Inches(0.4),
    'NestJS API', font_size=14, color=TEXT, bold=True, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(5.0), Inches(5.2), Inches(3), Inches(0.3),
    'Auth · Tests · Questions · Admin', font_size=11, color=TEXT_MUTED, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(9.0), Inches(4.9), Inches(3), Inches(0.4),
    'PostgreSQL / Supabase', font_size=14, color=TEXT, bold=True, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(9.0), Inches(5.2), Inches(3), Inches(0.3),
    'Prisma ORM · 6 Models', font_size=11, color=TEXT_MUTED, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(1.2), Inches(5.8), Inches(11), Inches(0.3),
    'API calls via Axios (JWT)    →    Report via UUID Token    →    Brevo Email    →    Supabase Storage', font_size=12, color=TEXT_MUTED, alignment=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════
# SLIDE 5 — 4 ELEMENTS & PERSONAS
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_accent_line(slide, Inches(0.8), Inches(0.8), Inches(2), GOLD)
add_text_box(slide, Inches(0.8), Inches(1.0), Inches(11), Inches(0.8),
    'The 4 Elements & Personas', font_size=38, color=WHITE, bold=True)

elements = [
    ('🔥 API (Fire)', 'Choleric', 'Berani, Intens & Bersemangat', 'Pemimpin alami, energik, tegas', API_RED, '#b83232'),
    ('💧 AIR (Water)', 'Melancholic', 'Dalam, Puitis & Introspektif', 'Pemikir mendalam, menghargai detail', AIR_BLUE, '#3b77bc'),
    ('🍃 ANGIN (Wind)', 'Sanguine', 'Cerah, Segar & Menginspirasi', 'Jiwa sosial, membawa keceriaan', ANGIN_GOLD, '#c5963a'),
    ('🌿 TANAH (Earth)', 'Phlegmatic', 'Tenang, Stabil & Harmonis', 'Sumber ketenangan, bijaksana', TANAH_GREEN, '#3a5c3a'),
]

for i, (name, parfum, tagline, desc, color, hex_c) in enumerate(elements):
    x = 0.8 + (i % 4) * 3.1
    y = 2.2
    card = add_shape_bg(slide, Inches(x), Inches(y), Inches(2.7), Inches(4.8), RGBColor(0x1c, 0x1b, 0x19))
    card.line.color.rgb = color
    card.line.width = Pt(1.5)
    add_text_box(slide, Inches(x + 0.2), Inches(y + 0.3), Inches(2.3), Inches(0.5),
        name, font_size=20, color=WHITE, bold=True)
    add_text_box(slide, Inches(x + 0.2), Inches(y + 0.9), Inches(2.3), Inches(0.4),
        f'Parfum: {parfum}', font_size=14, color=color, bold=True)
    add_text_box(slide, Inches(x + 0.2), Inches(y + 1.5), Inches(2.3), Inches(0.4),
        f'"{tagline}"', font_size=12, color=GOLD)
    add_text_box(slide, Inches(x + 0.2), Inches(y + 2.2), Inches(2.3), Inches(2.0),
        desc, font_size=13, color=RGBColor(0xa0, 0x98, 0x90))

# ═══════════════════════════════════════════════════════════════
# SLIDE 6 — USER JOURNEY
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_WARM)
add_accent_line(slide, Inches(0.8), Inches(0.8), Inches(2), GOLD)
add_text_box(slide, Inches(0.8), Inches(1.0), Inches(11), Inches(0.8),
    'User Journey', font_size=38, color=TEXT, bold=True)

steps = [
    ('1', 'Homepage', 'Hero canvas, 4 stats cards, CTA, element teasers'),
    ('2', 'Auth', 'Login / Register via email or Google OAuth'),
    ('3', 'Test (25 Questions)', '6-6-6-6 distribution + 1 closing survey; animated card stack'),
    ('4', 'Teaser / Reveal', 'Countdown animation → persona reveal with 3D wireframe background'),
    ('5', 'Report Page', 'Radar chart, persona cards, Shopee/TikTok links, feedback, share'),
]

for i, (num, title, desc) in enumerate(steps):
    y = 2.0 + i * 1.0
    add_shape_bg(slide, Inches(0.8), Inches(y), Inches(0.6), Inches(0.6), GOLD)
    add_text_box(slide, Inches(0.8), Inches(y + 0.05), Inches(0.6), Inches(0.5),
        num, font_size=22, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(1.8), Inches(y), Inches(3), Inches(0.4),
        title, font_size=18, color=TEXT, bold=True)
    add_text_box(slide, Inches(1.8), Inches(y + 0.35), Inches(10), Inches(0.4),
        desc, font_size=13, color=TEXT_MUTED)

# Arrow connecting steps
for i in range(4):
    y = 2.0 + i * 1.0 + 0.8
    add_text_box(slide, Inches(1.1), Inches(y), Inches(0.5), Inches(0.3),
        '↓', font_size=14, color=GOLD, bold=True, alignment=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════
# SLIDE 7 — FRONTEND PAGES
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_accent_line(slide, Inches(0.8), Inches(0.8), Inches(2), GOLD)
add_text_box(slide, Inches(0.8), Inches(1.0), Inches(11), Inches(0.8),
    'Frontend Pages', font_size=38, color=WHITE, bold=True)

pages = [
    ('/', 'HomePage', 'Canvas hero, 4 stat counters, CTA button, element teaser cards, donation badge'),
    ('/auth', 'AuthPage', 'Login/Register tabs, Google OAuth button, form validation via react-hook-form'),
    ('/test', 'TestPage', 'Protected route, animated card stack, progress bar, 25 randomized questions'),
    ('/teaser', 'TeaserPage', 'Countdown animation, wireframe 3D background, persona reveal'),
    ('/report/:token', 'ReportPage', 'Radar chart, persona primer/sekunder, shop links, share (html2canvas), feedback'),
    ('/profile', 'ProfilePage', 'User info, last persona, test history with links to previous reports'),
    ('/admin/*', 'AdminPanel', 'Dashboard stats, CRUD questions, paginated results, report templates editor'),
]
y = 2.2
for route, comp, desc in pages:
    add_text_box(slide, Inches(0.8), Inches(y), Inches(2.5), Inches(0.35),
        route, font_size=14, color=GOLD, bold=True)
    add_text_box(slide, Inches(3.5), Inches(y), Inches(2.2), Inches(0.35),
        comp, font_size=14, color=WHITE)
    add_text_box(slide, Inches(5.8), Inches(y), Inches(6.5), Inches(0.35),
        desc, font_size=12, color=RGBColor(0xa0, 0x98, 0x90))
    y += 0.65

# ═══════════════════════════════════════════════════════════════
# SLIDE 8 — BACKEND API
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_WARM)
add_accent_line(slide, Inches(0.8), Inches(0.8), Inches(2), GOLD)
add_text_box(slide, Inches(0.8), Inches(1.0), Inches(11), Inches(0.8),
    'Backend API Endpoints', font_size=38, color=TEXT, bold=True)

endpoints = [
    ('Public', 'GET /api/metrics/public, POST /api/metrics/ping'),
    ('Auth', 'POST /api/auth/register, POST /api/auth/login, GET /api/auth/me, GET /api/auth/google'),
    ('Questions', 'GET /api/questions/session (25 random), CRUD + toggle (admin)'),
    ('Tests', 'POST /api/tests/submit, GET /api/tests/report/:token, GET /api/tests/history'),
    ('Admin', 'GET /api/admin/stats, PATCH /api/admin/sales, GET/PUT templates, GET results'),
]
y = 2.2
for module, endpoints_list in endpoints:
    add_text_box(slide, Inches(0.8), Inches(y), Inches(2.5), Inches(0.35),
        module, font_size=16, color=GOLD_DARK, bold=True)
    add_text_box(slide, Inches(3.5), Inches(y), Inches(9), Inches(0.35),
        endpoints_list, font_size=13, color=TEXT)
    y += 0.55

add_text_box(slide, Inches(0.8), Inches(y + 0.3), Inches(11), Inches(0.8),
    'Backend modules: auth, users, questions, tests, metrics, email, admin, storage, reports — NestJS modular architecture',
    font_size=14, color=TEXT_MUTED)

# ═══════════════════════════════════════════════════════════════
# SLIDE 9 — SCORING SYSTEM
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_accent_line(slide, Inches(0.8), Inches(0.8), Inches(2), GOLD)
add_text_box(slide, Inches(0.8), Inches(1.0), Inches(11), Inches(0.8),
    'Scoring & Question System', font_size=38, color=WHITE, bold=True)

scoring_items = [
    '25 questions per session: 6 × API + 6 × AIR + 6 × ANGIN + 6 × TANAH + 1 closing survey',
    'Each option awards points to one of the 4 elements (targetType)',
    'Highest-scoring element = Persona Primer, second = Persona Sekunder',
    'Tie-breaker: API > AIR > ANGIN > TANAH (priority order)',
    'Each test generates a UUID token (30-day expiry) for public report access',
    'Email sent via Brevo with personalized report link after submission',
    'Users can optionally provide feedback (rating 1-5 + text) on their report',
]
add_bullet_slide(slide, Inches(0.8), Inches(2.2), Inches(11), Inches(4.5), scoring_items, 16, RGBColor(0xe8, 0xe3, 0xda))

# ═══════════════════════════════════════════════════════════════
# SLIDE 10 — DATABASE SCHEMA
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_WARM)
add_accent_line(slide, Inches(0.8), Inches(0.8), Inches(2), GOLD)
add_text_box(slide, Inches(0.8), Inches(1.0), Inches(11), Inches(0.8),
    'Database Schema (6 Models)', font_size=38, color=TEXT, bold=True)

db_items = [
    'User — id, email, name, googleId, password (bcrypt), isAdmin, timestamps',
    'Question — id, text, element (API/AIR/ANGIN/TANAH), isActive, createdAt',
    'AnswerOption — id, questionId, text, targetType (element), order',
    'TestResult — id, userId, scores (4 elements), personaPrimer/Sekunder, reportToken (UUID, 30d), feedback',
    'SystemMetric — singleton row (id=1), manualSalesCount, totalVisitors',
    'ReportTemplate — id (ElementType), parfumName, descriptions, backgroundImage, shop links',
]
add_bullet_slide(slide, Inches(0.8), Inches(2.2), Inches(11), Inches(4.5), db_items, 15, TEXT)

# ═══════════════════════════════════════════════════════════════
# SLIDE 11 — DESIGN SYSTEM
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_accent_line(slide, Inches(0.8), Inches(0.8), Inches(2), GOLD)
add_text_box(slide, Inches(0.8), Inches(1.0), Inches(11), Inches(0.8),
    'Design System — "Quiet Luxury"', font_size=38, color=WHITE, bold=True)

design_items = [
    'Inspiration: vundiego.com — warm neutral palette with gold accents',
    'Background: #f8f6f2 (light) / #111110 (dark)',
    'Primary accent: #c9a96e (Gold) — CTA buttons & highlights',
    'Typography: Playfair Display (headings, serif) + Inter (body, sans-serif)',
    'Pixel art icons for each element (Pixelify Sans font)',
    '4 element colors: Red (#b83232), Blue (#3b77bc), Gold (#c5963a), Green (#3a5c3a)',
    'Framer Motion for card stack animations, page transitions, countdown',
    'CSS custom properties throughout — centralized design tokens',
]
add_bullet_slide(slide, Inches(0.8), Inches(2.2), Inches(11), Inches(4.5), design_items, 16, RGBColor(0xe8, 0xe3, 0xda))

# ═══════════════════════════════════════════════════════════════
# SLIDE 12 — ADMIN PANEL
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_WARM)
add_accent_line(slide, Inches(0.8), Inches(0.8), Inches(2), GOLD)
add_text_box(slide, Inches(0.8), Inches(1.0), Inches(11), Inches(0.8),
    'Admin Panel Features', font_size=38, color=TEXT, bold=True)

admin_items = [
    'Dashboard: Total tests, visitors, sales, element distribution chart',
    'Questions: Full CRUD — create, edit, delete, toggle active/inactive',
    'Results: Paginated table of all test results with filters',
    'Templates: Edit report content per element (parfum name, descriptions, links)',
    'Sales counter: Manually update sales count for homepage display',
    'Protected by JWT + admin guard — only users with isAdmin=true can access',
    'Responsive sidebar layout with drawer for mobile',
]
add_bullet_slide(slide, Inches(0.8), Inches(2.2), Inches(11), Inches(4.5), admin_items, 16, TEXT)

# ═══════════════════════════════════════════════════════════════
# SLIDE 13 — KEY FEATURES
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_accent_line(slide, Inches(0.8), Inches(0.8), Inches(2), GOLD)
add_text_box(slide, Inches(0.8), Inches(1.0), Inches(11), Inches(0.8),
    'Key Features & Highlights', font_size=38, color=WHITE, bold=True)

features = [
    ('Dual Auth', 'Email/password + Google OAuth with account merging'),
    ('Auto-scored', 'Balanced 6-6-6-6 question distribution + tie-breaker logic'),
    ('Token Sharing', '30-day UUID tokens — share reports without login'),
    ('Visitor Metrics', 'SessionStorage de-duplication, 5-min cached stats'),
    ('Email Integration', 'Brevo API sends HTML report emails post-test'),
    ('Social Share', 'html2canvas generates PNG share cards client-side'),
    ('Admin CRUD', 'Full question bank & report template management'),
    ('Modern UI', 'Framer Motion animations, Canvas backgrounds, radar charts'),
    ('Security', 'bcrypt(12), JWT guards, admin guards, input validation'),
]
y = 2.2
for feat, desc in features:
    add_text_box(slide, Inches(0.8), Inches(y), Inches(2.5), Inches(0.35),
        feat, font_size=15, color=GOLD, bold=True)
    add_text_box(slide, Inches(3.5), Inches(y), Inches(9), Inches(0.35),
        desc, font_size=14, color=RGBColor(0xa0, 0x98, 0x90))
    y += 0.5

# ═══════════════════════════════════════════════════════════════
# SLIDE 14 — IMPLEMENTATION STATUS
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_WARM)
add_accent_line(slide, Inches(0.8), Inches(0.8), Inches(2), GOLD)
add_text_box(slide, Inches(0.8), Inches(1.0), Inches(11), Inches(0.8),
    'Implementation Status', font_size=38, color=TEXT, bold=True)

status_items = [
    'Fase 0: Project Setup (Vite, NestJS, Prisma, Supabase) ✅',
    'Fase 1: Design System (CSS tokens, UI components, Framer Motion) ✅',
    'Fase 2: Auth System (Register, Login, JWT, Google OAuth) ✅',
    'Fase 3: Homepage (Hero canvas, stats, CTA, element cards) ✅',
    'Fase 4: Test Interface (Card stack, 25 questions, progress) ✅',
    'Fase 5: Scoring & Report (Logic, radar chart, share, email) ✅',
    'Fase 6: Admin Panel (Dashboard, CRUD, results, templates) ✅',
    'Fase 7: Polish (Animations, SEO metadata, error handling) ✅',
    'Fase 8: Deployment (Vercel, Railway, CI/CD) ⏳ Pending',
]
add_bullet_slide(slide, Inches(0.8), Inches(2.2), Inches(11), Inches(4.5), status_items, 15, TEXT)

# ═══════════════════════════════════════════════════════════════
# SLIDE 15 — THANK YOU
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_shape_bg(slide, Inches(0), Inches(0), Inches(0.15), Inches(7.5), GOLD)
add_text_box(slide, Inches(1.5), Inches(2.5), Inches(10), Inches(1.2),
    'Terima Kasih', font_size=60, color=GOLD, bold=True, alignment=PP_ALIGN.LEFT)
add_accent_line(slide, Inches(1.5), Inches(3.8), Inches(3), GOLD)
add_text_box(slide, Inches(1.5), Inches(4.3), Inches(10), Inches(0.8),
    '4Persona Vundiego — Quiet Luxury', font_size=26, color=WHITE)
add_text_box(slide, Inches(1.5), Inches(5.0), Inches(10), Inches(0.5),
    'React 19 · NestJS 11 · PostgreSQL · Supabase · Prisma 6', font_size=16, color=RGBColor(0xa0, 0x98, 0x90))

# ── Save ──
output_path = 'C:\\Users\\isyaa\\OneDrive\\Documents\\Web and Code\\4Persona\\4Persona_Vundiego_Overview.pptx'
prs.save(output_path)
print(f'Presentation saved to: {output_path}')
print(f'Total slides: {len(prs.slides)}')
